from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
import re
import requests
from io import BytesIO
from PIL import Image

def parse_markdown(markdown):
    """
    Parses the given markdown text into a list of slides, where each slide is a tuple
    containing the title and a list of content items. Each content item is a tuple
    indicating its type ('bullet', 'text', 'image', or 'link') and the content string.
    """
    slides = markdown.split('---')
    parsed_slides = []
    for slide in slides:
        lines = slide.strip().split('\n')
        title = lines[0].replace('## ', '').strip()
        content = []
        for line in lines[1:]:
            if line.startswith('* '):
                content.append(('bullet', line[2:].strip()))
            elif line.startswith('!['):  # Image
                match = re.search(r'!\[.*\]\((.*)\)', line)
                if match:
                    content.append(('image', match.group(1).strip()))
            elif line.startswith('['):  # Hyperlink
                match = re.search(r'\[(.*)\]\((.*)\)', line)
                if match:
                    content.append(('link', match.groups()))
            else:
                content.append(('text', line.strip()))
        parsed_slides.append((title, content))
    return parsed_slides

def add_run_with_formatting(p, text, bold=False, italic=False, underline=False, strikethrough=False, highlight=False, color=None, hyperlink=None):
    """
    Adds a text run to the paragraph with the specified formatting options.
    """
    run = p.add_run()
    run.text = text
    if bold:
        run.font.bold = True
    if italic:
        run.font.italic = True
    if underline:
        run.font.underline = True
    if highlight:
        run.font.highlight = True  # Note: This feature may not be directly supported.
    if color:
        run.font.color.rgb = RGBColor.from_string(color)
    if hyperlink:
        run.hyperlink.address = hyperlink
    if strikethrough:
        # Workaround for strikethrough
        rPr = run._r.get_or_add_rPr()
        strike = OxmlElement('a:strike')
        strike.set('val', 'sngStrike')
        rPr.append(strike)

def apply_formatting(p, text):
    """
    Applies formatting to the given text within a paragraph based on markdown-like
    syntax for bold, italic, underline, strikethrough, highlight, color, and hyperlink.
    """
    parts = re.split(r'(\*\*|__|~~|\*|==|<color:#[0-9a-fA-F]{6}>|</color>|\[.*\]\(.*\))', text)
    format_stack = []
    color_stack = []
    hyperlink = None

    for part in parts:
        if part in ('**', '__'):
            if format_stack and format_stack[-1] == 'bold':
                format_stack.pop()
            else:
                format_stack.append('bold')
        elif part == '~~':
            if format_stack and format_stack[-1] == 'strikethrough':
                format_stack.pop()
            else:
                format_stack.append('strikethrough')
        elif part == '*':
            if format_stack and format_stack[-1] == 'italic':
                format_stack.pop()
            else:
                format_stack.append('italic')
        elif part == '==':
            if format_stack and format_stack[-1] == 'highlight':
                format_stack.pop()
            else:
                format_stack.append('highlight')
        elif part.startswith('<color:#'):
            color = part[7:13]
            color_stack.append(color)
        elif part == '</color>':
            if color_stack:
                color_stack.pop()
        elif re.match(r'\[.*\]\(.*\)', part):
            match = re.search(r'\[(.*)\]\((.*)\)', part)
            if match:
                text = match.group(1)
                hyperlink = match.group(2)
                add_run_with_formatting(p, text, hyperlink=hyperlink)
                hyperlink = None
        else:
            bold = 'bold' in format_stack
            italic = 'italic' in format_stack
            underline = 'underline' in format_stack
            strikethrough = 'strikethrough' in format_stack
            highlight = 'highlight' in format_stack
            color = color_stack[-1] if color_stack else None
            if color:
                color = color.lstrip('#')
            add_run_with_formatting(p, part, bold, italic, underline, strikethrough, highlight, color)

def fetch_image(url):
    """
    Fetches the image from the given URL and returns a BytesIO object if successful.
    """
    try:
        response = requests.get(url)
        response.raise_for_status()
        content_type = response.headers.get('Content-Type')
        if 'image' in content_type:
            return BytesIO(response.content)
        else:
            print(f"URL does not point to an image: {url}")
    except requests.exceptions.RequestException as e:
        print(f"Failed to fetch image from URL: {url}, error: {e}")
    return None

def create_ppt(parsed_slides, output_file):
    """
    Creates a PowerPoint presentation from the parsed slides and saves it to the given
    output file.
    """
    prs = Presentation()

    for title, content in parsed_slides:
        slide_layout = prs.slide_layouts[1] # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_box = slide.shapes.title
        content_box = slide.placeholders[1]
        
        title_box.text = title
        
        # Clear any existing text in the content box
        content_box.text_frame.clear()

        for item_type, item_content in content:
            if item_type == 'image':
                image_stream = fetch_image(item_content)
                if image_stream:
                    # Place the image on the right side
                    slide.shapes.add_picture(image_stream, Inches(6.5), Inches(1), width=Inches(2.5))
            else:
                if content_box.text_frame.text:  # Check if there is existing text
                    p = content_box.text_frame.add_paragraph()
                else:
                    p = content_box.text_frame.paragraphs[0]

                if item_type == 'text':
                    apply_formatting(p, item_content)
                elif item_type == 'bullet':
                    p.text = item_content
                    p.level = 1
                elif item_type == 'link':
                    text, url = item_content
                    add_run_with_formatting(p, text, hyperlink=url)
    
    prs.save(output_file)