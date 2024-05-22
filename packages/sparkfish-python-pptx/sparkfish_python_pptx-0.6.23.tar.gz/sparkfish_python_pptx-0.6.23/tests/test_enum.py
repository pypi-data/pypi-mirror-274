# encoding: utf-8

"""Unit-test suite for `pptx.enum` subpackage, focused on base classes.

Configured a little differently because of the meta-programming, the two enumeration
classes at the top constitute the entire fixture and most of the tests themselves just
make assertions on those.
"""

import pytest

from pptx.enum.base import (
    alias,
    Enumeration,
    EnumMember,
    ReturnValueOnlyEnumMember,
    XmlEnumeration,
    XmlMappedEnumMember,
)
from pptx.enum.shapes import PROG_ID, _ProgIdEnum
from pptx.util import Emu


@alias("BARFOO")
class FOOBAR(Enumeration):
    """
    Enumeration docstring
    """

    __ms_name__ = "MsoFoobar"

    __url__ = "http://msdn.microsoft.com/foobar.aspx"

    __members__ = (
        EnumMember("READ_WRITE", 1, "Readable and settable"),
        ReturnValueOnlyEnumMember("READ_ONLY", -2, "Return value only"),
    )


@alias("XML-FU")
class XMLFOO(XmlEnumeration):
    """
    XmlEnumeration docstring
    """

    __ms_name__ = "MsoXmlFoobar"

    __url__ = "http://msdn.microsoft.com/msoxmlfoobar.aspx"

    __members__ = (
        XmlMappedEnumMember(None, None, None, "No setting"),
        XmlMappedEnumMember("XML_RW", 42, "attrVal", "Read/write setting"),
        ReturnValueOnlyEnumMember("RO", -2, "Return value only;"),
    )


class DescribeEnumeration(object):
    def it_has_the_right_metaclass(self):
        assert type(FOOBAR).__name__ == "MetaEnumeration"

    def it_provides_an_EnumValue_instance_for_each_named_member(self):
        for obj in (FOOBAR.READ_WRITE, FOOBAR.READ_ONLY):
            assert type(obj).__name__ == "EnumValue"

    def it_provides_the_enumeration_value_for_each_named_member(self):
        assert FOOBAR.READ_WRITE == 1
        assert FOOBAR.READ_ONLY == -2

    def it_knows_if_a_setting_is_valid(self):
        FOOBAR.validate(FOOBAR.READ_WRITE)
        with pytest.raises(ValueError):
            FOOBAR.validate("foobar")
        with pytest.raises(ValueError):
            FOOBAR.validate(FOOBAR.READ_ONLY)

    def it_can_be_referred_to_by_a_convenience_alias_if_defined(self):
        assert BARFOO is FOOBAR  # noqa


class DescribeEnumValue(object):
    def it_provides_its_symbolic_name_as_its_string_value(self):
        assert ("%s" % FOOBAR.READ_WRITE) == "READ_WRITE (1)"

    def it_provides_its_description_as_its_docstring(self):
        assert FOOBAR.READ_ONLY.__doc__ == "Return value only"


class DescribeXmlEnumeration(object):
    def it_knows_the_XML_value_for_each_of_its_xml_members(self):
        assert XMLFOO.to_xml(XMLFOO.XML_RW) == "attrVal"
        assert XMLFOO.to_xml(42) == "attrVal"
        with pytest.raises(ValueError):
            XMLFOO.to_xml(XMLFOO.RO)

    def it_can_map_each_of_its_xml_members_from_the_XML_value(self):
        assert XMLFOO.from_xml(None) is None
        assert XMLFOO.from_xml("attrVal") == XMLFOO.XML_RW
        assert str(XMLFOO.from_xml("attrVal")) == "XML_RW (42)"


class Describe_ProgIdEnum(object):
    """Unit-test suite for `pptx.enum.shapes._ProgIdEnum."""

    def it_provides_access_to_its_members(self):
        assert type(PROG_ID.XLSX) == _ProgIdEnum.Member

    def it_can_test_an_item_for_membership(self):
        assert PROG_ID.XLSX in PROG_ID

    def it_has_a_readable_representation_for_itself(self):
        assert repr(PROG_ID) == "pptx.enum.shapes.PROG_ID"

    def it_has_a_readable_representation_for_each_of_its_members(self):
        assert repr(PROG_ID.XLSX) == "PROG_ID.XLSX"

    def it_has_attributes_on_each_member(self):
        assert PROG_ID.XLSX.height == Emu(609600)
        assert PROG_ID.XLSX.icon_filename == "xlsx-icon.emf"
        assert PROG_ID.XLSX.progId == "Excel.Sheet.12"
        assert PROG_ID.XLSX.width == Emu(965200)
